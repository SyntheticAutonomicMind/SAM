#!/usr/bin/env python3
# SPDX-License-Identifier: GPL-3.0-only
# SPDX-FileCopyrightText: Copyright (c) 2025 Andrew Wyatt (Fewtarius)

"""
Generate PowerPoint presentations from markdown using python-pptx.
Called by SAM's PPTXGenerator.swift

Input: JSON via stdin
Output: JSON with success status and file path

Usage:
    echo '{"slides": [...], "output": "output.pptx"}' | python3 generate_pptx.py
"""

import sys
import json
import logging
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    logging.error(f"Failed to import python-pptx: {e}")
    logging.error("Install with: pip install python-pptx")
    sys.exit(1)


def create_presentation(data):
    """
    Create PPTX from structured data.
    
    Args:
        data: {
            "slides": [
                {
                    "title": "Slide Title",
                    "content": ["Bullet 1", "Bullet 2"],  # Optional
                    "image": "path/to/image.png"          # Optional
                },
                ...
            ],
            "template": "path/to/template.pptx" (optional),
            "output": "path/to/output.pptx"
        }
        
    Returns:
        Path to generated PPTX file
    """
    # Load template or create blank presentation
    if data.get("template") and Path(data["template"]).exists():
        logging.info(f"Loading template: {data['template']}")
        prs = Presentation(data["template"])
    else:
        logging.info("Creating blank presentation")
        prs = Presentation()
    
    slides_data = data.get("slides", [])
    logging.info(f"Processing {len(slides_data)} slides")
    
    for i, slide_data in enumerate(slides_data, 1):
        logging.debug(f"Processing slide {i}: {slide_data.get('title', 'Untitled')}")
        
        # Determine slide type
        has_image = "image" in slide_data and slide_data["image"]
        has_content = "content" in slide_data and slide_data["content"]
        
        if has_image:
            # Image slide - use blank layout
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title text box (optional)
            if slide_data.get("title"):
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.2),
                    width=Inches(9), height=Inches(0.8)
                )
                title_frame = title_box.text_frame
                title_frame.text = slide_data["title"]
                title_frame.paragraphs[0].font.size = Pt(28)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add image
            image_path = slide_data["image"]
            if Path(image_path).exists():
                try:
                    # Center image on slide
                    slide.shapes.add_picture(
                        image_path,
                        Inches(1), Inches(1.5),
                        width=Inches(8)
                    )
                    logging.debug(f"Added image: {image_path}")
                except Exception as e:
                    logging.error(f"Failed to add image {image_path}: {e}")
            else:
                logging.warning(f"Image not found: {image_path}")
        
        elif has_content:
            # Content slide - use title and content layout
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data.get("title", "")
            
            # Add bullet points
            content_box = slide.placeholders[1]
            text_frame = content_box.text_frame
            text_frame.clear()  # Clear default text
            
            for j, line in enumerate(slide_data["content"]):
                if j == 0:
                    # First paragraph (already exists)
                    p = text_frame.paragraphs[0]
                else:
                    # Add new paragraph
                    p = text_frame.add_paragraph()
                
                p.text = line
                p.level = 0  # Top-level bullet
                p.font.size = Pt(18)
        
        else:
            # Title-only slide
            slide_layout = prs.slide_layouts[5]  # Title Only
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = slide_data.get("title", "")
    
    # Save presentation
    output_path = data["output"]
    prs.save(output_path)
    logging.info(f"Saved presentation: {output_path}")
    
    return output_path


def main():
    """Main entry point - read JSON from stdin, generate PPTX, output result."""
    # Configure logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(levelname)s: %(message)s'
    )
    
    try:
        # Read input JSON from stdin
        logging.debug("Reading input from stdin...")
        input_data = json.load(sys.stdin)
        
        # Validate required fields
        if "slides" not in input_data:
            raise ValueError("Missing required field: 'slides'")
        if "output" not in input_data:
            raise ValueError("Missing required field: 'output'")
        
        # Generate presentation
        output_path = create_presentation(input_data)
        
        # Output success result
        result = {
            "success": True,
            "path": output_path
        }
        print(json.dumps(result))
        sys.exit(0)
        
    except Exception as e:
        logging.error(f"Error generating presentation: {e}", exc_info=True)
        result = {
            "success": False,
            "error": str(e)
        }
        print(json.dumps(result))
        sys.exit(1)


if __name__ == "__main__":
    main()
